import cv2
import numpy as np
from PIL import Image
import os
from pptx import Presentation
from pptx.util import Inches
import tempfile
import shutil
from skimage.metrics import structural_similarity as ssim
import yt_dlp
import re
from urllib.parse import urlparse

class VideoToPPTConverter:
    def __init__(self, similarity_threshold=0.95, min_frame_interval=30):
        """
        Initialize the converter
        
        Args:
            similarity_threshold: Threshold for frame similarity (0-1, higher = more similar)
            min_frame_interval: Minimum frames between captures to avoid duplicates
        """
        self.similarity_threshold = similarity_threshold
        self.min_frame_interval = min_frame_interval
        
    def is_youtube_url(self, url):
        """
        Check if the provided string is a YouTube URL
        
        Args:
            url: String to check
            
        Returns:
            Boolean indicating if it's a YouTube URL
        """
        youtube_patterns = [
            r'(?:https?://)?(?:www\.)?youtube\.com/watch\?v=[\w-]+',
            r'(?:https?://)?(?:www\.)?youtu\.be/[\w-]+',
            r'(?:https?://)?(?:www\.)?youtube\.com/embed/[\w-]+',
            r'(?:https?://)?(?:www\.)?youtube\.com/v/[\w-]+'
        ]
        
        for pattern in youtube_patterns:
            if re.match(pattern, url):
                return True
        return False
    
    def download_youtube_video(self, youtube_url, output_dir="temp_downloads"):
        """
        Download YouTube video using yt-dlp
        
        Args:
            youtube_url: YouTube video URL
            output_dir: Directory to save downloaded video
            
        Returns:
            Path to downloaded video file
        """
        os.makedirs(output_dir, exist_ok=True)
        
        # Configure yt-dlp options
        ydl_opts = {
            'format': 'best[height<=720]',  # Download max 720p to save space
            'outtmpl': os.path.join(output_dir, '%(title)s.%(ext)s'),
            'quiet': False,
            'no_warnings': False,
        }
        
        print(f"Downloading video from: {youtube_url}")
        
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                # Extract video info
                info = ydl.extract_info(youtube_url, download=False)
                video_title = info.get('title', 'video')
                duration = info.get('duration', 0)
                
                print(f"Video: {video_title}")
                print(f"Duration: {duration//60}:{duration%60:02d}")
                
                # Download the video
                ydl.download([youtube_url])
                
                # Find the downloaded file
                for file in os.listdir(output_dir):
                    if file.startswith(video_title.replace('/', '_')):
                        return os.path.join(output_dir, file)
                
                # If exact match not found, return the most recent file
                files = [os.path.join(output_dir, f) for f in os.listdir(output_dir)]
                if files:
                    return max(files, key=os.path.getctime)
                
        except Exception as e:
            print(f"Error downloading video: {e}")
            raise Exception(f"Failed to download YouTube video: {e}")
        
    def extract_video_id(self, youtube_url):
        """
        Extract video ID from YouTube URL for filename
        
        Args:
            youtube_url: YouTube URL
            
        Returns:
            Video ID string
        """
        patterns = [
            r'(?:v=|\/)([0-9A-Za-z_-]{11}).*',
            r'(?:embed\/)([0-9A-Za-z_-]{11})',
            r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, youtube_url)
            if match:
                return match.group(1)
        
        return "unknown_video"
        
    def extract_key_frames(self, video_path, output_dir="temp_frames"):
        """
        Extract key frames when significant changes occur
        
        Args:
            video_path: Path to the input video file
            output_dir: Directory to save extracted frames
            
        Returns:
            List of frame file paths
        """
        # Create output directory
        os.makedirs(output_dir, exist_ok=True)
        
        # Open video
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            raise Exception(f"Error opening video file: {video_path}")
            
        frame_paths = []
        prev_frame = None
        frame_count = 0
        saved_count = 0
        last_saved_frame = -self.min_frame_interval
        
        print("Extracting key frames...")
        
        while True:
            ret, frame = cap.read()
            if not ret:
                break
                
            # Convert to grayscale for comparison
            gray_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            
            # Check if this is the first frame or significantly different
            is_key_frame = False
            
            if prev_frame is None:
                is_key_frame = True
            elif frame_count - last_saved_frame >= self.min_frame_interval:
                # Calculate similarity with previous saved frame
                similarity = ssim(prev_frame, gray_frame)
                if similarity < self.similarity_threshold:
                    is_key_frame = True
            
            if is_key_frame:
                # Save frame
                frame_filename = f"frame_{saved_count:04d}.png"
                frame_path = os.path.join(output_dir, frame_filename)
                cv2.imwrite(frame_path, frame)
                frame_paths.append(frame_path)
                
                prev_frame = gray_frame.copy()
                last_saved_frame = frame_count
                saved_count += 1
                
                print(f"Saved frame {saved_count} at {frame_count}/{int(cap.get(cv2.CAP_PROP_FRAME_COUNT))}")
            
            frame_count += 1
            
            # Optional: Show progress for long videos
            if frame_count % 100 == 0:
                total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                progress = (frame_count / total_frames) * 100
                print(f"Progress: {progress:.1f}% ({frame_count}/{total_frames})")
        
        cap.release()
        print(f"Extracted {len(frame_paths)} key frames")
        return frame_paths
    
    def create_presentation(self, frame_paths, output_ppt="video_presentation.pptx"):
        """
        Create PowerPoint presentation from extracted frames
        
        Args:
            frame_paths: List of frame image paths
            output_ppt: Output PowerPoint file path
        """
        print("Creating PowerPoint presentation...")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        for i, frame_path in enumerate(frame_paths):
            # Add slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image to slide
            img = Image.open(frame_path)
            
            # Calculate dimensions to fit slide while maintaining aspect ratio
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            img_width, img_height = img.size
            img_ratio = img_width / img_height
            slide_ratio = slide_width / slide_height
            
            if img_ratio > slide_ratio:
                # Image is wider than slide ratio
                width = slide_width
                height = slide_width / img_ratio
                left = 0
                top = (slide_height - height) / 2
            else:
                # Image is taller than slide ratio
                height = slide_height
                width = slide_height * img_ratio
                left = (slide_width - width) / 2
                top = 0
            
            # Add image to slide
            slide.shapes.add_picture(frame_path, left, top, width, height)
            
            print(f"Added slide {i+1}/{len(frame_paths)}")
        
        # Save presentation
        prs.save(output_ppt)
        print(f"Presentation saved as: {output_ppt}")
        
        return output_ppt
    
    def process_video(self, video_input, output_ppt=None, cleanup_temp=True):
        """
        Complete process: download (if URL), extract frames and create PPT
        
        Args:
            video_input: YouTube URL or path to local video file
            output_ppt: Path for output PowerPoint (if None, saves in script directory)
            cleanup_temp: Whether to delete temporary files
        """
        # Set default output path to script directory if not specified
        if output_ppt is None:
            script_dir = os.path.dirname(os.path.abspath(__file__))
            
            # Create a meaningful filename based on input
            if self.is_youtube_url(video_input):
                # Extract video ID from YouTube URL for filename
                video_id = self.extract_video_id(video_input)
                output_ppt = os.path.join(script_dir, f"youtube_{video_id}_slides.pptx")
            else:
                # Use local video filename
                video_name = os.path.splitext(os.path.basename(video_input))[0]
                output_ppt = os.path.join(script_dir, f"{video_name}_slides.pptx")
        
        temp_dir = tempfile.mkdtemp(prefix="video_frames_")
        download_dir = None
        video_path = video_input
        
        try:
            # Check if input is YouTube URL
            if self.is_youtube_url(video_input):
                download_dir = tempfile.mkdtemp(prefix="youtube_download_")
                video_path = self.download_youtube_video(video_input, download_dir)
                print(f"Downloaded to: {video_path}")
            elif not os.path.exists(video_input):
                raise Exception(f"Video file not found: {video_input}")
            
            # Extract key frames
            frame_paths = self.extract_key_frames(video_path, temp_dir)
            
            if not frame_paths:
                raise Exception("No frames were extracted from the video")
            
            # Create presentation
            self.create_presentation(frame_paths, output_ppt)
            
            print(f"Successfully created presentation with {len(frame_paths)} slides")
            print(f"Saved to: {os.path.abspath(output_ppt)}")
            
            return output_ppt
            
        finally:
            # Cleanup temporary files
            if cleanup_temp:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
                    print("Cleaned up temporary frame files")
                    
                if download_dir and os.path.exists(download_dir):
                    shutil.rmtree(download_dir)
                    print("Cleaned up downloaded video files")

# Usage example and configuration options
def main():
    """
    Main function with usage examples
    """
    # Initialize converter with custom settings
    converter = VideoToPPTConverter(
        similarity_threshold=0.90,  # Adjust sensitivity (lower = more slides)
        min_frame_interval=30       # Minimum frames between captures (higher = fewer slides)
    )
    
    # Example inputs - you can use either YouTube URLs or local files
    examples = [
        "https://www.youtube.com/watch?v=pFze1OkCx30",  # YouTube URL
        "https://youtu.be/VIDEO_ID",                     # Short YouTube URL
        "local_video.mp4"                               # Local file path
    ]
    
    # Choose your input
    video_input = "https://www.youtube.com/watch?v=8_frhkUWErs"  # Replace with your URL/path
    # output_ppt will be automatically set to script directory if not specified
    
    try:
        print("Starting video processing...")
        print(f"Input: {video_input}")
        
        # Process video - output_ppt=None means save in script directory
        output_file = converter.process_video(video_input, output_ppt=None)
        print(f"Success! Presentation saved as: {output_file}")
        
        # Additional processing suggestions
        print("\nNext steps:")
        print(f"1. Open {os.path.basename(output_file)} to review the slides")
        print("2. Use the OCR extractor to get text from slides")
        print("3. Run: python ppt_ocr_extractor.py")
        
    except Exception as e:
        print(f"Error: {e}")
        print("\nTroubleshooting tips:")
        print("1. Check your internet connection for YouTube videos")
        print("2. Make sure the YouTube URL is valid and accessible")
        print("3. For local files, ensure the file exists and is readable")
        print("4. Install required packages: pip install yt-dlp")
        print("5. Try adjusting similarity_threshold (0.85-0.98)")

# Additional utility functions
def download_youtube_only(youtube_url, output_dir="downloads"):
    """
    Utility function to just download a YouTube video
    """
    converter = VideoToPPTConverter()
    return converter.download_youtube_video(youtube_url, output_dir)

def process_multiple_videos(video_list, output_dir=None):
    """
    Process multiple videos (URLs or files) into separate presentations
    """
    # If no output directory specified, use script directory
    if output_dir is None:
        output_dir = os.path.dirname(os.path.abspath(__file__))
    
    os.makedirs(output_dir, exist_ok=True)
    converter = VideoToPPTConverter(similarity_threshold=0.90, min_frame_interval=30)
    
    results = []
    for i, video_input in enumerate(video_list):
        try:
            output_name = f"presentation_{i+1:02d}.pptx"
            output_path = os.path.join(output_dir, output_name)
            
            print(f"\nProcessing {i+1}/{len(video_list)}: {video_input}")
            converter.process_video(video_input, output_path)
            
            results.append({
                'input': video_input,
                'output': output_path,
                'status': 'success'
            })
            
        except Exception as e:
            print(f"Failed to process {video_input}: {e}")
            results.append({
                'input': video_input,
                'output': None,
                'status': f'failed: {e}'
            })
    
    return results

if __name__ == "__main__":
    main()

# Installation requirements:
"""
Required packages:
pip install opencv-python python-pptx pillow scikit-image yt-dlp numpy

Additional notes:
- yt-dlp is the modern replacement for youtube-dl
- For older systems, you might need: pip install youtube-dl (less reliable)
- Make sure you have sufficient disk space for video downloads
- Some YouTube videos may be restricted by region or privacy settings

Usage examples:
1. YouTube URL: "https://www.youtube.com/watch?v=VIDEO_ID"
2. Short URL: "https://youtu.be/VIDEO_ID"  
3. Local file: "/path/to/video.mp4"
4. Playlist: Use process_multiple_videos() with list of URLs
"""